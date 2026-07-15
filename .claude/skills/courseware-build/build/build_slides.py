#!/usr/bin/env python3
"""Generate the "Build a Human-AI Workforce with Autonomous AI Agents" course slide deck (all-white Tertiary house style).

Design helpers are the same set used by the tertiary-course-slides skill that
produced the n8n reference deck (cover, section, content, two_col, cards3,
big_statement, step_slide, test_slide, brk). Content is driven entirely by
course_data.py + data_domainN.py so the deck stays 100% aligned with the LP,
LG and labs.
"""
import os, sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
import course_data as C
from data_domain1 import DOMAIN1
from data_domain2 import DOMAIN2
from data_domain3 import DOMAIN3
from data_domain4 import DOMAIN4
from data_domain5 import DOMAIN5
ACTIVITIES = DOMAIN1 + DOMAIN2 + DOMAIN3 + DOMAIN4 + DOMAIN5

def _find_repo(start):
    """Locate the course repo (a dir containing both courseware/ and labs/).
    Env COURSE_REPO overrides. Keeps the build working wherever the skill lives."""
    env = os.environ.get("COURSE_REPO")
    if env and os.path.isdir(env):
        return env
    d = start
    for _ in range(8):
        d = os.path.dirname(d)
        if os.path.isdir(os.path.join(d, "courseware")) and os.path.isdir(os.path.join(d, "labs")):
            return d
    return os.path.dirname(os.path.dirname(HERE))
REPO = _find_repo(HERE)
ASSETS = os.path.join(os.path.dirname(HERE), "assets")   # co-located with the skill

# ---------------- palette (matches reference) ----------------
BLUE=RGBColor(0x1F,0x6F,0xEB); TEAL=RGBColor(0x10,0xB9,0x81); AMBER=RGBColor(0xF5,0x9E,0x0B)
INK=RGBColor(0x16,0x1B,0x26); GREY=RGBColor(0x5B,0x63,0x72); LIGHT=RGBColor(0xF5,0xF8,0xFC)
WHITE=RGBColor(0xFF,0xFF,0xFF); LINE=RGBColor(0xE2,0xE8,0xF0); VIOLET=RGBColor(0x7C,0x3A,0xED)
ORANGE=RGBColor(0xEA,0x58,0x0C); RED=RGBColor(0xDC,0x26,0x26)
# Per-topic accent themes for the section covers: Hermes=orange, OpenClaw=red, Paperclip=blue.
TOPIC_THEME={1:ORANGE, 2:RED, 3:BLUE}

# ---------------- dark "masterclass" palette (Topic 1 / Hermes only) ----------------
DK_BG   = RGBColor(0x0D,0x0E,0x0B)   # slide background
DK_PANEL= RGBColor(0x15,0x17,0x12)   # cards/rows/panels
DK_AMBER= RGBColor(0xF5,0xA6,0x23)   # headlines, kickers, accents, number chips
DK_IVORY= RGBColor(0xE8,0xE6,0xDD)   # body text
DK_FAINT= RGBColor(0x2A,0x2C,0x28)   # divider rules
DK_DIM  = RGBColor(0x8A,0x8D,0x86)   # captions/footers
DK_GREEN= RGBColor(0x4A,0xDE,0x80)   # success/test-it accents

# Module-level theme switch — dark "masterclass" theme for the ENTIRE deck.
THEME={"dark":False}
def _bg():   return DK_BG    if THEME["dark"] else WHITE
def _panel():return DK_PANEL if THEME["dark"] else LIGHT
def _ink():  return DK_IVORY if THEME["dark"] else INK
def _grey(): return DK_DIM   if THEME["dark"] else GREY
def _line(): return DK_FAINT if THEME["dark"] else LINE
def _acc(c): return DK_AMBER if THEME["dark"] else c
THEME["dark"]=True   # whole deck renders in the dark masterclass theme

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height
BLANK=prs.slide_layouts[6]

def slide(): return prs.slides.add_slide(BLANK)
def rect(s,x,y,w,h,color,line=None):
    sp=s.shapes.add_shape(1,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1)
    sp.shadow.inherit=False; return sp
def oval(s,x,y,w,h,color):
    sp=s.shapes.add_shape(9,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    sp.line.fill.background(); sp.shadow.inherit=False; return sp
def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space=4):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(space)
        for t,sz,col,bold in line:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=bold
            r.font.color.rgb=col; r.font.name="Arial"
    return tb
def bullets(s,x,y,w,h,items,size=18,color=None,gap=10,mcolor=BLUE):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    body=color if color is not None else _ink()
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(gap)
        lvl=it[1] if isinstance(it,tuple) else 0
        text=it[0] if isinstance(it,tuple) else it
        r=p.add_run(); r.text=("•  " if lvl==0 else "–  ")+text
        r.font.size=Pt(size if lvl==0 else size-2); r.font.color.rgb=body if lvl==0 else _grey()
        r.font.name="Arial"; r.font.bold=(lvl==0 and isinstance(it,tuple) and len(it)>2 and it[2])
    return tb

PAGE={"n":0}
def footer(s):
    PAGE["n"]+=1
    txt(s,Inches(0.4),Inches(7.05),Inches(7.5),Inches(0.35),
        [[(f"{C.SHORT_TITLE}  ·  {C.COURSE_CODE}",9,_grey(),False)]])
    txt(s,Inches(5.0),Inches(7.05),Inches(3.3),Inches(0.35),
        [[("© 2026 Tertiary Infotech Academy Pte Ltd",9,_grey(),False)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(12.4),Inches(7.05),Inches(0.6),Inches(0.35),
        [[(str(PAGE["n"]),9,_grey(),False)]],align=PP_ALIGN.RIGHT)
def head(s,title,kicker=None,kcolor=BLUE):
    rect(s,0,0,SW,SH,_bg()); rect(s,0,0,Inches(0.28),Inches(1.55),_acc(kcolor))
    if kicker: txt(s,Inches(0.85),Inches(0.48),Inches(11.9),Inches(0.4),[[(kicker,14,_acc(kcolor),True)]])
    # Auto-fit the title so long lab titles never wrap into the divider rule below.
    L=len(str(title)); tsize = 29 if L<=42 else (25 if L<=56 else 21)
    txt(s,Inches(0.85),Inches(0.9),Inches(11.95),Inches(0.82),[[(title,tsize,_ink(),True)]],anchor=MSO_ANCHOR.MIDDLE)
    rect(s,Inches(0.85),Inches(1.8),Inches(11.63),Inches(0.02),_line())
    return s
def _logo(name):
    p=os.path.join(ASSETS,name)
    return p if os.path.exists(p) else None

# ---------------- slide templates ----------------
def cover():
    s=slide(); rect(s,0,0,SW,SH,_bg())
    rect(s,0,0,SW,Inches(0.22),BLUE); rect(s,0,Inches(7.28),SW,Inches(0.22),TEAL)
    org=_logo("tertiary-infotech-logo.png")
    if org: s.shapes.add_picture(org,Inches(0.85),Inches(0.7),height=Inches(1.05))
    # course badge (top-right) — use a neutral WSQ/course badge asset if one exists; otherwise no badge
    badge=_logo("wsq-badge.png") or _logo("course-badge.png")
    if badge:
        s.shapes.add_picture(badge,Inches(10.05),Inches(0.6),width=Inches(2.5))
    txt(s,Inches(0.9),Inches(2.3),Inches(12),Inches(0.6),[[("COURSE SLIDES  ·  WSQ",16,_acc(BLUE),True)]])
    txt(s,Inches(0.9),Inches(2.85),Inches(12.0),Inches(1.9),[[(C.TITLE,40,_ink(),True)]])
    rect(s,Inches(0.92),Inches(4.75),Inches(2.4),Inches(0.06),TEAL)
    txt(s,Inches(0.9),Inches(5.05),Inches(12),Inches(1.4),
        [[(f"WSQ Course Code: {C.COURSE_CODE}",16,_grey(),False)],
         [("Conducted by Tertiary Infotech Academy Pte Ltd  ·  UEN 201200696W",14,_grey(),False)]],space=6)
    txt(s,Inches(0.9),Inches(6.5),Inches(12),Inches(0.4),[[(f"Version {C.VERSION}  ·  {C.VERSION_DATE}",12,_grey(),False)]])
    txt(s,Inches(0.9),Inches(6.85),Inches(12),Inches(0.34),[[("© 2026 Tertiary Infotech Academy Pte Ltd. All rights reserved.  ·  www.tertiarycourses.com.sg",10,_grey(),False)]])

def section(kicker,title,n,sub="",accent=BLUE):
    s=slide(); rect(s,0,0,SW,SH,_bg()); rect(s,0,0,Inches(0.28),SH,accent)
    rect(s,Inches(0.85),Inches(2.5),Inches(0.14),Inches(2.0),accent)
    txt(s,Inches(1.25),Inches(2.55),Inches(11),Inches(0.6),[[(kicker,18,accent,True)]])
    txt(s,Inches(1.25),Inches(3.0),Inches(11.4),Inches(1.6),[[(title,40,_ink(),True)]])
    if sub: txt(s,Inches(1.27),Inches(4.55),Inches(11),Inches(0.8),[[(sub,16,_grey(),False)]])
    ghost=DK_FAINT if THEME["dark"] else RGBColor(0xE2,0xE8,0xF0)
    txt(s,Inches(10.0),Inches(0.7),Inches(2.8),Inches(1.6),[[(n,72,ghost,True)]],align=PP_ALIGN.RIGHT)
    footer(s)
def content(title,items,kicker=None,size=20):
    # HARD RULE: no bullet-wall slides — every content slide renders as numbered
    # visual cards (LIGHT card, accent bar, number chip), never plain bullets.
    import math as _m
    s=head(slide(),title,kicker)
    xs=[(i[0] if isinstance(i,tuple) else i) for i in items]
    n=max(1,len(xs)); cols=2 if n>3 else 1; rows=_m.ceil(n/cols)
    x0,y0=Inches(0.85),Inches(2.0); gw=Inches(11.63); gap=Inches(0.18)
    cw=int((gw-gap*(cols-1))/cols); ch=int((Inches(4.7)-gap*(rows-1))/rows)
    fs=(16 if rows<=2 else 14 if rows==3 else 12) if cols==2 else (17 if rows<=3 else 14)
    for idx,it in enumerate(xs):
        r,c=divmod(idx,cols)
        cx=x0+c*(cw+gap); cy=y0+r*(ch+gap)
        rect(s,cx,cy,cw,ch,_panel()); rect(s,cx,cy,Inches(0.07),ch,_acc(BLUE))
        oval(s,cx+Inches(0.2),int(cy+ch/2-Inches(0.19)),Inches(0.38),Inches(0.38),_acc(BLUE))
        txt(s,cx+Inches(0.2),int(cy+ch/2-Inches(0.17)),Inches(0.38),Inches(0.34),
            [[(str(idx+1),13,DK_BG if THEME["dark"] else WHITE,True)]],align=PP_ALIGN.CENTER)
        txt(s,cx+Inches(0.74),cy+Inches(0.08),cw-Inches(0.98),ch-Inches(0.16),
            [[(it,fs,_ink(),False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def two_col(title,left,right,kicker=None,lhead="",rhead=""):
    s=head(slide(),title,kicker)
    rect(s,Inches(0.85),Inches(1.95),Inches(5.7),Inches(4.7),_panel()); rect(s,Inches(6.95),Inches(1.95),Inches(5.55),Inches(4.7),_panel())
    if lhead: txt(s,Inches(1.1),Inches(2.15),Inches(5.2),Inches(0.4),[[(lhead,16,_acc(BLUE),True)]])
    if rhead: txt(s,Inches(7.2),Inches(2.15),Inches(5.0),Inches(0.4),[[(rhead,16,_acc(TEAL),True)]])
    bullets(s,Inches(1.1),Inches(2.7),Inches(5.2),Inches(3.8),left,size=16,color=_ink())
    bullets(s,Inches(7.2),Inches(2.7),Inches(5.05),Inches(3.8),right,size=16,color=_ink(),mcolor=TEAL); footer(s); return s
def cards3(title,cards,kicker):
    s=head(slide(),title,kicker); xs=[Inches(0.85),Inches(5.0),Inches(9.15)]
    for i,c in enumerate(cards[:3]):
        x=xs[i]; col=c[0]
        rect(s,x,Inches(1.95),Inches(3.65),Inches(4.7),_panel()); rect(s,x,Inches(1.95),Inches(3.65),Inches(0.12),col)
        txt(s,x+Inches(0.25),Inches(2.2),Inches(3.2),Inches(0.6),[[(c[1],19,col,True)]])
        bullets(s,x+Inches(0.25),Inches(2.95),Inches(3.2),Inches(3.4),c[2],size=14,color=_ink(),mcolor=col,gap=9)
    footer(s); return s
def big_statement(line1,line2,kicker,color=BLUE):
    s=slide(); rect(s,0,0,SW,SH,_bg()); rect(s,0,0,Inches(0.28),SH,color)
    txt(s,Inches(1.1),Inches(2.2),Inches(11),Inches(0.5),[[(kicker,16,color,True)]])
    txt(s,Inches(1.1),Inches(2.8),Inches(11.3),Inches(2.4),[[(line1,38,_ink(),True)]])
    if line2: txt(s,Inches(1.12),Inches(4.9),Inches(11),Inches(1.2),[[(line2,20,_grey(),False)]])
    footer(s); return s
import math
PALETTE=[BLUE,TEAL,VIOLET,AMBER]
def tile_grid(title,items,kicker=None,cols=2,size=15,icons=None,accent=BLUE):
    """Grid of light panels, each with a coloured icon/number badge + text.
    items: list of strings (or (title,caption) tuples). Much richer than a bullet list."""
    s=head(slide(),title,kicker,kcolor=accent)
    n=len(items); rows=math.ceil(n/cols)
    X0=Inches(0.85); Y0=Inches(1.95); TOTW=Inches(11.63); AREAH=Inches(4.78)
    gx=Inches(0.3); gy=Inches(0.26)
    cw=int((TOTW-gx*(cols-1))/cols); ch=int((AREAH-gy*(rows-1))/rows)
    bd=Inches(0.6)
    for i,it in enumerate(items):
        r=i//cols; c=i%cols
        x=int(X0+(cw+gx)*c); y=int(Y0+(ch+gy)*r); col=PALETTE[i%len(PALETTE)]
        rect(s,x,y,cw,ch,_panel()); rect(s,x,y,Inches(0.1),ch,col)
        oval(s,x+Inches(0.28),int(y+ch/2-bd/2),bd,bd,col)
        ic=icons[i] if icons else str(i+1)
        txt(s,x+Inches(0.28),int(y+ch/2-bd/2),bd,bd,[[(ic,19,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        tx=x+Inches(1.08); tw=cw-Inches(1.32)
        if isinstance(it,tuple):
            txt(s,tx,int(y+Inches(0.14)),tw,int(ch-Inches(0.2)),
                [[(it[0],size+2,_ink(),True)],[(it[1],size-2,_grey(),False)]],anchor=MSO_ANCHOR.MIDDLE,space=3)
        else:
            txt(s,tx,int(y+Inches(0.1)),tw,int(ch-Inches(0.16)),[[(it,size,_ink(),False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def flow_h(title,steps,kicker=None,color=BLUE):
    """Horizontal numbered flow: coloured chips connected by chevrons."""
    s=head(slide(),title,kicker,kcolor=color)
    n=len(steps); X0=Inches(0.85); TOTW=Inches(11.63); gap=Inches(0.34)
    cw=int((TOTW-gap*(n-1))/n); y=Inches(2.55); ch=Inches(3.15); bd=Inches(0.82)
    for i,st in enumerate(steps):
        x=int(X0+(cw+gap)*i)
        rect(s,x,y,cw,ch,_panel()); rect(s,x,y,cw,Inches(0.1),color)
        oval(s,int(x+cw/2-bd/2),int(y+Inches(0.42)),bd,bd,color)
        txt(s,int(x+cw/2-bd/2),int(y+Inches(0.42)),bd,bd,[[(str(i+1),30,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,x+Inches(0.16),int(y+Inches(1.55)),cw-Inches(0.32),int(ch-Inches(1.7)),[[(st,14,_ink(),False)]],align=PP_ALIGN.CENTER)
        if i<n-1:
            txt(s,int(x+cw-Inches(0.04)),int(y+ch/2-Inches(0.3)),int(gap+Inches(0.08)),Inches(0.6),
                [[("▶",15,color,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def trainer_slide(kicker,name,role,rows,initials,accent=BLUE):
    """Profile-card layout: avatar badge + name/role panel on the left, labelled
    info tiles on the right. rows: list of (LABEL, value); blank value → fill-in line."""
    s=head(slide(),"About the Trainer",kicker,kcolor=accent)
    lx=Inches(0.85); lw=Inches(3.65)
    rect(s,lx,Inches(1.95),lw,Inches(4.7),_panel()); rect(s,lx,Inches(1.95),lw,Inches(0.12),accent)
    bd=Inches(1.7); ax=int(lx+(lw-bd)/2)
    oval(s,ax,Inches(2.5),bd,bd,accent)
    txt(s,ax,Inches(2.5),bd,bd,[[(initials,44,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,lx+Inches(0.15),Inches(4.55),lw-Inches(0.3),Inches(0.6),[[(name,21,_ink(),True)]],align=PP_ALIGN.CENTER)
    txt(s,lx+Inches(0.15),Inches(5.2),lw-Inches(0.3),Inches(1.2),[[(role,13,_grey(),False)]],align=PP_ALIGN.CENTER)
    rx=Inches(4.9); rw=Inches(7.6); ry=Inches(1.95); rh=Inches(4.7)
    n=len(rows); gy=Inches(0.2); th=int((rh-gy*(n-1))/n)
    for i,(label,val) in enumerate(rows):
        y=int(ry+(th+gy)*i); col=PALETTE[i%len(PALETTE)]
        rect(s,rx,y,rw,th,_panel()); rect(s,rx,y,Inches(0.1),th,col)
        vruns=[(val,14,_ink(),False)] if val else [("____________________________________________",13,_line(),False)]
        txt(s,rx+Inches(0.32),y,rw-Inches(0.6),th,
            [[(label.upper(),11,col,True)],vruns],anchor=MSO_ANCHOR.MIDDLE,space=3)
    footer(s); return s
def activity_overview(tag,title,desc,build,services,kicker):
    s=head(slide(),title,kicker,kcolor=TEAL)
    rect(s,Inches(0.85),Inches(1.85),Inches(1.7),Inches(0.5),_acc(TEAL))
    txt(s,Inches(0.85),Inches(1.9),Inches(1.7),Inches(0.4),[[(tag,16,DK_BG if THEME["dark"] else WHITE,True)]],align=PP_ALIGN.CENTER)
    # Auto-fit the description so long lab goals never clip into the panel below.
    dL=len(str(desc)); dsize = 21 if dL<=200 else (17 if dL<=320 else 15)
    txt(s,Inches(0.85),Inches(2.55),Inches(11.7),Inches(1.6),[[(desc,dsize,_ink(),False)]])
    rect(s,Inches(0.85),Inches(4.3),Inches(11.7),Inches(2.0),_panel())
    txt(s,Inches(1.1),Inches(4.5),Inches(11),Inches(0.4),[[("You'll build",14,_acc(BLUE),True)]])
    txt(s,Inches(1.1),Inches(4.9),Inches(11),Inches(0.6),[[(build,18,_ink(),True)]])
    txt(s,Inches(1.1),Inches(5.6),Inches(11.2),Inches(0.6),[[("Tools:  ",13,_grey(),True),(services,13,_grey(),False)]]); footer(s); return s
def step_slide(kicker,act_title,n,total,text,cmd="",prompt=""):
    s=head(slide(),act_title,kicker,TEAL)
    oval(s,Inches(0.85),Inches(2.5),Inches(1.4),Inches(1.4),_acc(TEAL))
    txt(s,Inches(0.85),Inches(2.74),Inches(1.4),Inches(0.9),[[(str(n),38,DK_BG if THEME["dark"] else WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.95),Inches(1.95),Inches(11),Inches(0.4),[[(f"STEP {n} OF {total}",13,_grey(),True)]])
    txt(s,Inches(2.55),Inches(2.4),Inches(10.1),Inches(1.3),[[(text,23,_ink(),False)]],anchor=MSO_ANCHOR.MIDDLE)
    # HARD RULE: never render a comment-only "command" (# …) on a slide — the
    # code box appears only for a real, runnable command.
    if cmd and not cmd.lstrip().startswith("#"):
        if THEME["dark"]:
            # slightly lighter fill + faint border so the box reads against the dark bg
            rect(s,Inches(2.55),Inches(4.15),Inches(10.1),Inches(0.95),RGBColor(0x10,0x14,0x1E),line=DK_FAINT)
        else:
            rect(s,Inches(2.55),Inches(4.15),Inches(10.1),Inches(0.95),RGBColor(0x0B,0x12,0x20))
        txt(s,Inches(2.8),Inches(4.28),Inches(9.7),Inches(0.7),[[("$ "+cmd,13,RGBColor(0x9C,0xDC,0xFE),False)]],anchor=MSO_ANCHOR.MIDDLE)
    if prompt:
        # the natural-language alternative — most users prompt the agent instead
        # of typing the CLI command, so every command step shows both
        py=Inches(5.32) if (cmd and not cmd.lstrip().startswith("#")) else Inches(4.15)
        rect(s,Inches(2.55),py,Inches(10.1),Inches(0.95),_panel())
        rect(s,Inches(2.55),py,Inches(0.05),Inches(0.95),_acc(TEAL))
        txt(s,Inches(2.8),py,Inches(9.75),Inches(0.95),
            [[("💬 Or just ask:  ",13.5,_acc(TEAL),True),("\u201c"+prompt+"\u201d",13.5,_ink(),False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def dark_rows(tag,title_lines,sub_lines,rows,warn=None,accent=None,numbered=True):
    """Masterclass-style dark feature slide: terminal tag, amber (or custom
    accent) headline, light subtitle lines, numbered/bulleted row-cards. Rows
    may be strings, (amber,white) tuples, or lists of (text,colorkey) pairs
    (keys: w/a/b/g/d). warn=(BADGE, text) renders a red warning chip."""
    DKBG=RGBColor(0x0D,0x0E,0x0B); AMBER=RGBColor(0xF5,0xA6,0x23)
    ROWBG=RGBColor(0x15,0x17,0x12); FAINT=RGBColor(0x2A,0x2C,0x28); IVORY=RGBColor(0xE8,0xE6,0xDD)
    WARNRED=RGBColor(0xE0,0x5A,0x5A)
    ACC=accent or AMBER
    s=slide(); rect(s,0,0,SW,SH,DKBG)
    txt(s,Inches(0.85),Inches(0.4),Inches(11.6),Inches(0.35),[[("> "+tag,13,ACC,True)]])
    rect(s,Inches(0.85),Inches(0.82),Inches(11.63),Inches(0.02),FAINT)
    y=0.98
    for tl in title_lines:
        txt(s,Inches(0.85),Inches(y),Inches(11.9),Inches(0.78),[[(tl.upper(),40,ACC,True)]]); y+=0.76
    y+=0.08
    for sl in sub_lines:
        txt(s,Inches(0.85),Inches(y),Inches(11.9),Inches(0.38),[[(sl,15,IVORY,False)]]); y+=0.4
    y+=0.22
    n=max(1,len(rows)); gap=0.14
    bottom=6.35 if warn else 6.85
    rh=min(0.78,(bottom-y-gap*(n-1))/n)
    for i,r in enumerate(rows):
        ry=Inches(y+i*(rh+gap))
        rect(s,Inches(0.85),ry,Inches(11.63),Inches(rh),ROWBG)
        rect(s,Inches(0.85),ry,Inches(0.05),Inches(rh),ACC)
        if numbered:
            txt(s,Inches(1.1),ry,Inches(0.75),Inches(rh),[[(f"{i+1:02d}",16,ACC,True)]],anchor=MSO_ANCHOR.MIDDLE)
            tx,tw=Inches(1.9),Inches(10.35)
        else:
            txt(s,Inches(1.1),ry,Inches(0.4),Inches(rh),[[("▸",14,ACC,True)]],anchor=MSO_ANCHOR.MIDDLE)
            tx,tw=Inches(1.55),Inches(10.7)
        txt(s,tx,ry,tw,Inches(rh),[_dkruns(r,16)],anchor=MSO_ANCHOR.MIDDLE)
    if warn:
        wy=y+n*(rh+gap)+0.08
        chipw=max(1.15,0.11*len(warn[0])+0.3)   # chip grows with the badge text so it never wraps
        rect(s,Inches(0.85),Inches(wy),Inches(chipw),Inches(0.34),WARNRED)
        txt(s,Inches(0.85),Inches(wy+0.02),Inches(chipw),Inches(0.3),[[(warn[0],12,DKBG,True)]],align=PP_ALIGN.CENTER)
        txt(s,Inches(1.05+chipw),Inches(wy),Inches(11.28-chipw),Inches(0.34),[[(warn[1],13,IVORY,False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
_DKC={"w":RGBColor(0xE8,0xE6,0xDD),"a":RGBColor(0xF5,0xA6,0x23),"b":RGBColor(0x6E,0xA8,0xFF),
      "g":RGBColor(0x4A,0xDE,0x80),"d":RGBColor(0x8A,0x8D,0x86),"r":RGBColor(0xE8,0x6A,0x8A)}
def _dkruns(spec,size=16):
    """spec: str | (amber,white) tuple | list of (text,colorkey) pairs -> txt() runs line."""
    if isinstance(spec,str): return [(spec,size,_DKC["w"],False)]
    if isinstance(spec,tuple): return [(spec[0],size,_DKC["a"],True),(spec[1],size,_DKC["w"],False)]
    return [(t,size,_DKC.get(k,_DKC["w"]),k in ("a","b","g","r")) for t,k in spec]
def dark_table(tag,title,col1,col2,rows,foot=None):
    """Masterclass-style dark comparison table: amber vs blue column headers,
    white row labels, thin divider rules, optional mixed-colour footer line."""
    DKBG=RGBColor(0x0D,0x0E,0x0B); AMBER=_DKC["a"]; HBLUE=_DKC["b"]
    FAINT=RGBColor(0x2A,0x2C,0x28); IVORY=_DKC["w"]; DIM=_DKC["d"]
    s=slide(); rect(s,0,0,SW,SH,DKBG)
    txt(s,Inches(0.85),Inches(0.4),Inches(8.0),Inches(0.35),[[("> "+tag,13,AMBER,True)]])
    txt(s,Inches(8.9),Inches(0.4),Inches(3.6),Inches(0.35),[[("● WHEN TO USE WHICH",12,DIM,True)]],align=PP_ALIGN.RIGHT)
    rect(s,Inches(0.85),Inches(0.82),Inches(11.63),Inches(0.02),FAINT)
    txt(s,Inches(0.85),Inches(1.0),Inches(11.9),Inches(0.8),[[(title.upper(),40,AMBER,True)]])
    LX,C1,C2=Inches(0.85),Inches(3.3),Inches(8.1)
    y=2.05
    txt(s,C1,Inches(y),Inches(4.6),Inches(0.4),[[(col1.upper(),15,AMBER,True)]])
    txt(s,C2,Inches(y),Inches(4.3),Inches(0.4),[[(col2.upper(),15,HBLUE,True)]])
    y+=0.5
    rh=0.62
    for label,v1,v2 in rows:
        rect(s,LX,Inches(y),Inches(11.63),Inches(0.015),FAINT)
        txt(s,LX,Inches(y+0.08),Inches(2.3),Inches(rh),[[(label,15,IVORY,True)]])
        txt(s,C1,Inches(y+0.08),Inches(4.6),Inches(rh),[[(v1,15,IVORY,False)]])
        txt(s,C2,Inches(y+0.08),Inches(4.3),Inches(rh),[[(v2,15,IVORY,False)]])
        y+=rh
    rect(s,LX,Inches(y),Inches(11.63),Inches(0.015),FAINT)
    if foot:
        txt(s,Inches(0.85),Inches(y+0.3),Inches(11.63),Inches(0.45),[_dkruns(foot,15)],align=PP_ALIGN.CENTER)
    footer(s); return s
def dark_code(tag,title,sub,code_lines,foot=None,win_title="",right_tag=None):
    """Masterclass-style dark slide with a terminal/code window: amber headline,
    window chrome (● ● ● + title), monospace coloured code lines, amber footer."""
    DKBG=RGBColor(0x0D,0x0E,0x0B); AMBER=_DKC["a"]; IVORY=_DKC["w"]; FAINT=RGBColor(0x2A,0x2C,0x28)
    CODEBG=RGBColor(0x08,0x0A,0x08); CHROME=RGBColor(0x14,0x16,0x12); DIM=_DKC["d"]
    s=slide(); rect(s,0,0,SW,SH,DKBG)
    txt(s,Inches(0.85),Inches(0.4),Inches(8.0),Inches(0.35),[[("> "+tag,13,AMBER,True)]])
    if right_tag:
        txt(s,Inches(8.9),Inches(0.4),Inches(3.6),Inches(0.35),[[("● "+right_tag,12,DIM,True)]],align=PP_ALIGN.RIGHT)
    rect(s,Inches(0.85),Inches(0.82),Inches(11.63),Inches(0.02),FAINT)
    txt(s,Inches(0.85),Inches(0.95),Inches(11.9),Inches(0.75),[[(title.upper(),38,AMBER,True)]])
    y=1.72
    if sub:
        txt(s,Inches(0.85),Inches(y),Inches(11.9),Inches(0.4),[[(sub,15,IVORY,False)]]); y+=0.46
    # window chrome
    ch=0.34
    rect(s,Inches(0.85),Inches(y),Inches(11.63),Inches(ch),CHROME)
    for k,cc in enumerate((RGBColor(0xE0,0x5A,0x5A),RGBColor(0xF5,0xA6,0x23),RGBColor(0x4A,0xDE,0x80))):
        oval(s,Inches(1.02+k*0.22),Inches(y+0.11),Inches(0.12),Inches(0.12),cc)
    if win_title:
        txt(s,Inches(1.8),Inches(y+0.02),Inches(9.5),Inches(0.3),[[(win_title,11,DIM,False)]])
    y+=ch
    n=len(code_lines); lh=min(0.34,(6.35-y)/max(1,n))
    rect(s,Inches(0.85),Inches(y),Inches(11.63),Inches(lh*n+0.2),CODEBG)
    for i,ln in enumerate(code_lines):
        runs=_dkruns(ln,13)
        # monospace look: Courier New via manual run construction
        tb=s.shapes.add_textbox(Inches(1.1),Inches(y+0.08+i*lh),Inches(11.1),Inches(lh))
        tf=tb.text_frame; tf.word_wrap=False
        p=tf.paragraphs[0]
        for t_,sz,col,bold in runs:
            r=p.add_run(); r.text=t_; r.font.size=Pt(12.5); r.font.bold=bold
            r.font.color.rgb=col; r.font.name="Courier New"
    if foot:
        txt(s,Inches(0.85),Inches(6.62),Inches(11.63),Inches(0.4),[[(foot,14,AMBER,True)]],align=PP_ALIGN.CENTER)
    footer(s); return s
def dark_catalog(tag,title,cards,scope=None,right_tag=None):
    """Masterclass-style dark slide: a row of amber-outlined cards, each with a
    heading + a list of tool names, plus an optional SCOPE note panel below."""
    DKBG=RGBColor(0x0D,0x0E,0x0B); AMBER=_DKC["a"]; IVORY=_DKC["w"]
    FAINT=RGBColor(0x2A,0x2C,0x28); CARDBG=RGBColor(0x10,0x12,0x0E); DIM=_DKC["d"]
    s=slide(); rect(s,0,0,SW,SH,DKBG)
    txt(s,Inches(0.85),Inches(0.4),Inches(8.0),Inches(0.35),[[("> "+tag,13,AMBER,True)]])
    if right_tag:
        txt(s,Inches(8.4),Inches(0.4),Inches(4.1),Inches(0.35),[[("● "+right_tag,12,DIM,True)]],align=PP_ALIGN.RIGHT)
    rect(s,Inches(0.85),Inches(0.82),Inches(11.63),Inches(0.02),FAINT)
    txt(s,Inches(0.85),Inches(0.98),Inches(11.9),Inches(0.78),[[(title.upper(),40,AMBER,True)]])
    y=1.95
    n=len(cards); gap=Inches(0.18); cw=int((Inches(11.63)-gap*(n-1))/n); chh=Inches(3.0); b=Inches(0.02)
    for i,(headt,items) in enumerate(cards):
        x=int(Inches(0.85)+i*(cw+gap)); yv=Inches(y)
        rect(s,x,yv,cw,chh,CARDBG)
        rect(s,x,yv,cw,b,AMBER); rect(s,x,int(yv+chh-b),cw,b,AMBER)
        rect(s,x,yv,b,chh,AMBER); rect(s,int(x+cw-b),yv,b,chh,AMBER)
        txt(s,x,int(yv+Inches(0.14)),cw,Inches(0.34),[[(headt.upper(),13,AMBER,True)]],align=PP_ALIGN.CENTER)
        rows=[[(it,12,IVORY if not it.startswith("+") and "·" not in it else DIM,False)] for it in items]
        txt(s,x+Inches(0.18),int(yv+Inches(0.58)),cw-Inches(0.3),int(chh-Inches(0.7)),rows,space=6)
    if scope:
        sy=y+3.0+0.25
        rect(s,Inches(0.85),Inches(sy),Inches(11.63),Inches(1.15),CARDBG)
        rect(s,Inches(0.85),Inches(sy),Inches(0.05),Inches(1.15),AMBER)
        txt(s,Inches(1.1),Inches(sy+0.12),Inches(1.3),Inches(0.35),[[("SCOPE",15,AMBER,True)]])
        txt(s,Inches(2.5),Inches(sy+0.1),Inches(9.7),Inches(0.98),[[(scope,12.5,IVORY,False)]])
    footer(s); return s
def dark_panels(tag,title,panels,foot=None,right_tag=None):
    """Masterclass-style dark slide: coloured outlined panels ('// HEAD' + body
    run-lines) and a mixed-colour footer. panels: (colorkey, head, [runlines])."""
    DKBG=RGBColor(0x0D,0x0E,0x0B); FAINT=RGBColor(0x2A,0x2C,0x28); PBG=RGBColor(0x10,0x12,0x0E)
    AMBER=_DKC["a"]; DIM=_DKC["d"]
    PCOL={"a":_DKC["a"],"b":_DKC["b"],"g":_DKC["g"],"r":RGBColor(0xE8,0x6A,0x8A)}
    s=slide(); rect(s,0,0,SW,SH,DKBG)
    txt(s,Inches(0.85),Inches(0.4),Inches(8.0),Inches(0.35),[[("> "+tag,13,AMBER,True)]])
    if right_tag:
        txt(s,Inches(8.4),Inches(0.4),Inches(4.1),Inches(0.35),[[("● "+right_tag,12,DIM,True)]],align=PP_ALIGN.RIGHT)
    rect(s,Inches(0.85),Inches(0.82),Inches(11.63),Inches(0.02),FAINT)
    txt(s,Inches(0.85),Inches(0.98),Inches(11.9),Inches(0.78),[[(title.upper(),40,AMBER,True)]])
    y=2.0; n=len(panels); gap=Inches(0.25)
    pw=int((Inches(11.63)-gap*(n-1))/n); ph=Inches(3.3); b=Inches(0.02)
    for i,(ck,head,body) in enumerate(panels):
        col=PCOL.get(ck,AMBER); x=int(Inches(0.85)+i*(pw+gap)); yv=Inches(y)
        rect(s,x,yv,pw,ph,PBG)
        rect(s,x,yv,pw,b,col); rect(s,x,int(yv+ph-b),pw,b,col)
        rect(s,x,yv,b,ph,col); rect(s,int(x+pw-b),yv,b,ph,col)
        txt(s,x+Inches(0.28),int(yv+Inches(0.2)),pw-Inches(0.5),Inches(0.4),[[("// "+head.upper(),15,col,True)]])
        txt(s,x+Inches(0.28),int(yv+Inches(0.72)),pw-Inches(0.52),int(ph-Inches(0.9)),
            [_dkruns(ln,13) for ln in body],space=6)
    if foot:
        txt(s,Inches(0.85),Inches(y+3.3+0.28),Inches(11.63),Inches(0.45),[_dkruns(foot,15)],align=PP_ALIGN.CENTER)
    footer(s); return s
def dark_flow(tag,title,steps,foot=None,right_tag=None):
    """Masterclass-style dark slide: amber-outlined step boxes joined by arrows,
    plus a mixed-colour footer. steps: (head, body)."""
    DKBG=RGBColor(0x0D,0x0E,0x0B); FAINT=RGBColor(0x2A,0x2C,0x28); PBG=RGBColor(0x10,0x12,0x0E)
    AMBER=_DKC["a"]; IVORY=_DKC["w"]; DIM=_DKC["d"]
    s=slide(); rect(s,0,0,SW,SH,DKBG)
    txt(s,Inches(0.85),Inches(0.4),Inches(8.0),Inches(0.35),[[("> "+tag,13,AMBER,True)]])
    if right_tag:
        txt(s,Inches(8.4),Inches(0.4),Inches(4.1),Inches(0.35),[[("● "+right_tag,12,DIM,True)]],align=PP_ALIGN.RIGHT)
    rect(s,Inches(0.85),Inches(0.82),Inches(11.63),Inches(0.02),FAINT)
    txt(s,Inches(0.85),Inches(1.05),Inches(11.9),Inches(0.85),[[(title.upper(),44,AMBER,True)]])
    y=2.55; n=len(steps); ag=Inches(0.38); b=Inches(0.02)
    bw=int((Inches(11.63)-ag*(n-1))/n); bh=Inches(1.55)
    for i,(head,body) in enumerate(steps):
        x=int(Inches(0.85)+i*(bw+ag)); yv=Inches(y)
        rect(s,x,yv,bw,bh,PBG)
        rect(s,x,yv,bw,b,AMBER); rect(s,x,int(yv+bh-b),bw,b,AMBER)
        rect(s,x,yv,b,bh,AMBER); rect(s,int(x+bw-b),yv,b,bh,AMBER)
        txt(s,x,int(yv+Inches(0.16)),bw,Inches(0.36),[[(head.upper(),14,AMBER,True)]],align=PP_ALIGN.CENTER)
        txt(s,x+Inches(0.12),int(yv+Inches(0.6)),bw-Inches(0.24),Inches(0.85),[[(body,11.5,IVORY,False)]],align=PP_ALIGN.CENTER)
        if i<n-1:
            txt(s,int(x+bw),int(yv+bh/2-Inches(0.2)),ag,Inches(0.4),[[("→",18,AMBER,True)]],align=PP_ALIGN.CENTER)
    if foot:
        txt(s,Inches(0.85),Inches(y+1.55+0.45),Inches(11.63),Inches(0.8),[_dkruns(foot,15)],align=PP_ALIGN.CENTER)
    footer(s); return s
def dark_cards(tag,title_lines,cards,notes=(),badge=None):
    """Masterclass-style dark slide: amber headline + a row of green-outlined cards
    (name + subtitle) + optional AUTO-INSTALLED badge line and dim footnotes."""
    DKBG=RGBColor(0x0D,0x0E,0x0B); AMBER=RGBColor(0xF5,0xA6,0x23); GREEN=RGBColor(0x4A,0xDE,0x80)
    IVORY=RGBColor(0xE8,0xE6,0xDD); FAINT=RGBColor(0x2A,0x2C,0x28); CARDBG=RGBColor(0x10,0x12,0x0E)
    DIM=RGBColor(0x8A,0x8D,0x86)
    s=slide(); rect(s,0,0,SW,SH,DKBG)
    txt(s,Inches(0.85),Inches(0.4),Inches(11.6),Inches(0.35),[[("> "+tag,13,AMBER,True)]])
    rect(s,Inches(0.85),Inches(0.82),Inches(11.63),Inches(0.02),FAINT)
    y=0.98
    for tl in title_lines:
        txt(s,Inches(0.85),Inches(y),Inches(11.9),Inches(0.78),[[(tl.upper(),40,AMBER,True)]]); y+=0.76
    y+=0.35
    n=len(cards); gap=Inches(0.25); cw=int((Inches(11.63)-gap*(n-1))/n); ch=Inches(1.55); b=Inches(0.025)
    for i,(name,sub) in enumerate(cards):
        x=int(Inches(0.85)+i*(cw+gap)); yv=Inches(y)
        rect(s,x,yv,cw,ch,CARDBG)
        rect(s,x,yv,cw,b,GREEN); rect(s,x,int(yv+ch-b),cw,b,GREEN)
        rect(s,x,yv,b,ch,GREEN); rect(s,int(x+cw-b),yv,b,ch,GREEN)
        txt(s,x,int(yv+Inches(0.2)),cw,Inches(0.62),[[(name,26,GREEN,True)]],align=PP_ALIGN.CENTER)
        txt(s,x,int(yv+Inches(0.92)),cw,Inches(0.42),[[(sub,15,IVORY,False)]],align=PP_ALIGN.CENTER)
    y+=1.55+0.4
    rest=list(notes)
    if badge and rest:
        rect(s,Inches(0.85),Inches(y),Inches(2.0),Inches(0.36),AMBER)
        txt(s,Inches(0.85),Inches(y+0.02),Inches(2.0),Inches(0.32),[[(badge,12,DKBG,True)]],align=PP_ALIGN.CENTER)
        txt(s,Inches(3.0),Inches(y),Inches(9.4),Inches(0.36),[[(rest.pop(0),14,IVORY,True)]],anchor=MSO_ANCHOR.MIDDLE)
        y+=0.5
    for nt in rest:
        txt(s,Inches(0.85),Inches(y),Inches(11.6),Inches(0.34),[[(nt,12,DIM,False)]]); y+=0.36
    footer(s); return s
def dark_pairs(tag,title_lines,sub,pairs):
    """Masterclass-style dark slide: amber headline + a 2-column grid of
    command/description cells (amber command, light description)."""
    import math as _m
    DKBG=RGBColor(0x0D,0x0E,0x0B); AMBER=RGBColor(0xF5,0xA6,0x23)
    IVORY=RGBColor(0xE8,0xE6,0xDD); FAINT=RGBColor(0x2A,0x2C,0x28); ROWBG=RGBColor(0x15,0x17,0x12)
    DIM=RGBColor(0x8A,0x8D,0x86)
    s=slide(); rect(s,0,0,SW,SH,DKBG)
    txt(s,Inches(0.85),Inches(0.4),Inches(11.6),Inches(0.35),[[("> "+tag,13,AMBER,True)]])
    rect(s,Inches(0.85),Inches(0.82),Inches(11.63),Inches(0.02),FAINT)
    y=0.98
    for tl in title_lines:
        txt(s,Inches(0.85),Inches(y),Inches(11.9),Inches(0.78),[[(tl.upper(),40,AMBER,True)]]); y+=0.76
    if sub:
        txt(s,Inches(0.85),Inches(y+0.02),Inches(11.9),Inches(0.36),[[(sub,14,DIM,False)]]); y+=0.44
    y+=0.18
    cols=2; rows=_m.ceil(len(pairs)/cols); gx=Inches(0.3); gy=0.13
    cw=int((Inches(11.63)-gx)/cols); rh=min(0.8,(6.8-y-gy*(rows-1))/rows)
    for i,(cmd,desc) in enumerate(pairs):
        r,c=i%rows,i//rows   # fill column-first like the reference
        x=int(Inches(0.85)+c*(cw+gx)); yv=Inches(y+r*(rh+gy))
        rect(s,x,yv,cw,Inches(rh),ROWBG)
        rect(s,x,yv,Inches(0.045),Inches(rh),AMBER)
        txt(s,x+Inches(0.22),yv,Inches(2.35),Inches(rh),[[(cmd,14,AMBER,True)]],anchor=MSO_ANCHOR.MIDDLE)
        txt(s,x+Inches(2.6),yv,cw-Inches(2.8),Inches(rh),[[(desc,12.5,IVORY,False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def shot(title,img,kicker=None,caption="",frame=True):
    """Screenshot slide — a real UI capture with a caption (any aspect).
    frame=False renders the image bare so dark, pre-blended figures merge
    with the deck background instead of looking pasted."""
    s=head(slide(),title,kicker,TEAL)
    try:
        from PIL import Image as _PILImage
        _iw,_ih=_PILImage.open(img).size
    except Exception:
        _iw,_ih=1440,900
    H=Inches(4.5); W=int(H*_iw/_ih)
    maxW=Inches(11.6)
    if W>maxW: W=int(maxW); H=int(maxW*_ih/_iw)
    x=int((SW-W)/2); y=int(Inches(2.02)+(Inches(4.5)-H)/2)
    if frame:
        rect(s,x-Inches(0.06),y-Inches(0.06),W+Inches(0.12),H+Inches(0.12),_line())
    s.shapes.add_picture(img,x,y,width=W,height=H)
    if caption:
        _csz=12 if len(caption)<=130 else 10.5   # shrink long captions so they never wrap onto the footer
        txt(s,Inches(0.85),Inches(6.68),Inches(11.6),Inches(0.35),[[(caption,_csz,_grey(),False)]],align=PP_ALIGN.CENTER)
    footer(s); return s
def test_slide(act_title,text,kicker):
    s=head(slide(),act_title,kicker,TEAL)
    rect(s,Inches(0.85),Inches(2.3),Inches(11.7),Inches(2.6),DK_PANEL if THEME["dark"] else RGBColor(0xE8,0xF7,0xEE))
    txt(s,Inches(1.2),Inches(2.6),Inches(11),Inches(0.5),[[("✅  Test it",20,DK_GREEN if THEME["dark"] else RGBColor(0x12,0x7A,0x3E),True)]])
    txt(s,Inches(1.2),Inches(3.3),Inches(11),Inches(1.4),[[(text,18,_ink(),False)]]); footer(s); return s
def brk(kind,dur,color=AMBER):
    s=slide(); rect(s,0,0,SW,SH,_bg())
    rect(s,0,0,SW,Inches(0.22),color); rect(s,0,Inches(7.28),SW,Inches(0.22),color)
    rect(s,Inches(5.4),Inches(2.35),Inches(2.53),Inches(0.1),color)
    txt(s,0,Inches(2.75),SW,Inches(1.2),[[(kind,48,_ink(),True)]],align=PP_ALIGN.CENTER)
    txt(s,0,Inches(4.05),SW,Inches(0.8),[[(dur,22,color,True)]],align=PP_ALIGN.CENTER); PAGE["n"]+=1

# ============================================================ BUILD
cover()

# ---------------- ADMIN ----------------
section("COURSE ADMINISTRATION","Welcome & Housekeeping","")
content("Digital Attendance (Mandatory)",[
 "It is mandatory to take the AM, PM and Assessment digital attendance for WSQ-funded courses.",
 "The trainer/administrator displays the digital attendance QR code from the SSG portal.",
 "Scan the QR code with your mobile phone camera and submit your attendance.",
 "A minimum of 75% attendance is required to be eligible for assessment and funding."],kicker="TRAQOM · SSG DIGITAL ATTENDANCE")
trainer_slide("YOUR TRAINER · GENERAL","Your Trainer","General Trainer template —\nto be completed by the trainer",
 [("Name",""),("Title / Designation",""),("Qualifications",""),
  ("Areas of expertise",""),("Training & industry experience",""),("Contact","")],
 initials="?",accent=GREY)
trainer_slide("YOUR TRAINER",C.TRAINER,"Principal Trainer\nTertiary Infotech Academy Pte. Ltd.",
 [("Role","Principal Trainer, Tertiary Infotech Academy Pte. Ltd."),
  ("Certification","AI, automation and cloud certified — building autonomous AI agents and agentic workflows."),
  ("Delivers","WSQ courses on AI agents, automation, data engineering and cloud."),
  ("Founder","Founder and lead instructor at Tertiary Infotech / Tertiary Courses.")],
 initials="AA",accent=BLUE)
content("Let's Know Each Other",[
 "Your name and organisation / role.",
 "Your experience with AI, automation or coding (if any).",
 "What you want to be able to automate with AI agents after this course."],kicker="ICE-BREAKER")
tile_grid("Ground Rules",[
 "Set your mobile phone to silent mode.","Participate actively — no question is too small.",
 "Mutual respect: agree to disagree.","One conversation at a time.",
 "Be punctual; return from breaks on time.","75% attendance is required."],
 kicker="HOUSEKEEPING",cols=2,size=15)
# Live screenshot of the LMS-TMS portal (house standard: a visual, not a text link)
shot("Download Course Material",
     os.path.join(REPO,"courseware","assets","screenshots","lms-portal.png"),
     kicker="COURSE PORTAL",
     caption="Sign in at https://lms-tms.tertiaryinfotech.com → open this course → Courseware tab → download the Slides, Learner Guide and Lesson Plan.")
# Lesson plan overview — rendered from the day themes so it can never drift
two_col("Lesson Plan — 2 Days, 8 hours/day",[
 (f"Day 1 — {C.DAY_THEMES[1]}",0),
 ("Topic 1: Hermes Agent (Labs 1–14)",1),
 ("Topic 2: OpenClaw (Labs 15–26)",1)],
 [(f"Day 2 — {C.DAY_THEMES[2]}",0),
 ("Topic 3: Paperclip (Labs 27–32)",1),
 ("Final Assessment (WA + PP)",1),
 ("Daily timing",0),
 ("9:30am–6:30pm · 1-hour lunch · tea breaks within",1)],
 kicker="SCHEDULE",lhead="Day 1",rhead="Day 2")
# Learning-outcome tiles built straight from course_data (strip the "LOn: " prefix)
_LO_TITLES=["Analyze AI Applications","Design & Chatbot Efficiency","Evaluate & Improve RAG"]
_lo_tiles=[(_LO_TITLES[i] if i<len(_LO_TITLES) else f"LO{i+1}",
            lo.split(": ",1)[1] if ": " in lo else lo)
           for i,lo in enumerate(C.LEARNING_OUTCOMES)]
tile_grid("Learning Outcomes",_lo_tiles,kicker="WHAT YOU'LL ACHIEVE",cols=2,size=15)
content("Briefing for Assessment",[
 "Place phones and other materials under the table or on the floor.",
 "No photos or recording of assessment scripts.","No discussion during the assessment.",
 "Use a black/blue pen for hard-copy assessments.","No liquid paper / correction tape.",
 "Scripts are collected when time is up."])
content("Assessment",[
 C.ASSESSMENT["written"], C.ASSESSMENT["practical"],
 "Format: Open Book — slides, Learner Guide and approved materials only.",
 C.ASSESSMENT["note"],"An appeal process is available if required."],kicker="FINAL ASSESSMENT")
flow_h("Assessment Flow",[
 "TRAQOM survey — scan the QR code on the LMS",
 "Assessment digital attendance — scan the SSG QR",
 "Sit WA (SAQ) then PP — open book",
 "Submit your answers on the LMS",
 "Sign the Assessment Summary Record"],kicker="ON ASSESSMENT DAY")

# ---------------- CORE AI-AGENT CONCEPTS ----------------
section("CORE CONCEPTS","Autonomous AI Agents Fundamentals","")
tile_grid("What is an Autonomous AI Agent?",[
 ("Perceive → reason → act → remember","An LLM-driven loop: it takes in a request, decides what to do, acts, then remembers the outcome."),
 ("Uses tools & skills","It calls web search, images, email, code and integrations — not just chat — to get real work done."),
 ("Runs across channels","The same agent is reachable on Telegram, WhatsApp, Discord, Slack and the terminal."),
 ("Autonomy under oversight","It works on its own, but important actions pause for human approval.")],
 kicker="OVERVIEW",cols=2,size=15)
tile_grid("The Agent Stack",[
 ("Model / LLM","The reasoning engine — Claude, OpenAI, OpenRouter, MiniMax or DeepSeek."),
 ("Memory","Durable, user- or company-scoped recall across sessions and surfaces."),
 ("Skills & Tools","Reusable capabilities plus tool and third-party integration calls."),
 ("Channels","Where the agent meets people: Telegram, WhatsApp, Discord, Slack."),
 ("Scheduler & Security","Crons and a heartbeat automate work; sandboxes, budgets and approvals contain it.")],
 kicker="THE BUILDING BLOCKS",cols=2,size=14)
two_col("RAG vs a Standard LLM Assistant",
 [("Retrieval Augmented Generation (RAG)",0),
  ("Retrieves from a curated, current corpus before generating",1),
  ("Grounded, auditable answers — less hallucination",1),
  ("Effectiveness hinges on corpus freshness, indexing and how retrieval/memory is configured",1)],
 [("Standard LLM assistant",0),
  ("Answers from training data alone — fastest to deploy",1),
  ("No live company data, no auditable response provenance",1),
  ("Evaluate both on accuracy, relevance, latency and maintenance cost",1)],
 kicker="PERFORMANCE EFFECTIVENESS",lhead="RAG agent",rhead="LLM-only assistant")
big_statement("Powerful, but always under human oversight.","Approvals gate risky actions, budgets cap spend (warn at 80%, pause at 100%), and sandboxes isolate execution — so an agent stays safe as it scales.","HUMAN-IN-THE-LOOP GOVERNANCE",color=BLUE)
_pf=cards3("Three Platforms, One Skillset",[
 (BLUE,"Hermes Agent",["Your personal chief-of-staff","Lives across your messaging apps","Cross-channel, memory, crons"]),
 (TEAL,"OpenClaw",["Automates a business back-office","Node.js gateway daemon","Skills, tools, nightly reports"]),
 (VIOLET,"Paperclip",["A company of AI agents you govern","Self-hosted with Docker","Hiring, budgets, audit trail"])],kicker="WHAT YOU'LL BUILD")
# The three platform logos, one on each card — on a white plate so every mark stays legible
_LOGODIR=os.path.join(REPO,"courseware","assets","logos")
for _i,_f in enumerate(["hermes-agent.png","openclaw.png","paperclip.png"]):
    _p=os.path.join(_LOGODIR,_f)
    if os.path.exists(_p):
        _cx=[Inches(0.85),Inches(5.0),Inches(9.15)][_i]
        _plate=Inches(1.7); _px=int(_cx+(Inches(3.65)-_plate)/2); _py=int(Inches(4.75))
        rect(_pf,_px,_py,int(_plate),int(_plate),WHITE,line=_line())
        _lg=Inches(1.4)
        _pf.shapes.add_picture(_p,int(_px+(_plate-_lg)/2),int(_py+(_plate-_lg)/2),width=int(_lg),height=int(_lg))
content("The Build Arc",[
 "Install the runtime → connect a model → add channels → give it memory.",
 "Extend it with skills → wire in tools and integrations → schedule work with crons.",
 "Secure it with sandboxing, secrets and approvals → drive it through commands.",
 "Finish each platform by orchestrating a multi-agent team for an end-to-end outcome.",
 "The same arc repeats on Hermes, OpenClaw and Paperclip — so the skillset transfers."],kicker="HOW EVERY LAB PROGRESSES")

# ---------------- TOPICS + ACTIVITIES ----------------
# Real screenshots captured from a live local Hermes install (dashboard on
# http://127.0.0.1:9119), plus the live skills-category snapshot — rendered as
# framed screenshot slides and category tile grids inside the matching labs.
import json as _json
_SHOTDIR=os.path.join(REPO,"courseware","assets","screenshots")
_SHOTS={2:[("hermes-dashboard.png","Dashboard","The Hermes Dashboard (http://127.0.0.1:9119) — your local deployment's sessions view"),
           ("hermes-channels.png","Channels","The Channels page (127.0.0.1:9119/channels) — 31 channels from one gateway; credentials in ~/.hermes/.env"),
           ("hermes-achievements.png","Achievements","The Achievements page — Hermes gamifies agent mastery with tiers (Copper → Silver → Gold → Diamond → Olympian) and unlockable badges")],
        4:[("hermes-skills.png","Skills","The Skills page of a live Hermes install — 97 skills, filterable by category, plus Learn-a-skill and New-skill")],
        5:[("hermes-env.png","Providers & API Keys","The Keys page (127.0.0.1:9119/env) — OAuth provider logins (Nous Portal, OpenAI, MiniMax, Anthropic, …) and API keys, stored in ~/.hermes/.env")],
        9:[("hermes-profiles.png","Profiles","The Profiles page (127.0.0.1:9119/profiles) — per-profile toolsets and defaults that shape what the agent may use"),
           ("hermes-kanban.png","Kanban","The Kanban board — supervise the agent's tasks as they move across the columns")]}
# Live screenshots from the running Paperclip company (http://127.0.0.1:3100/TER)
_PCDIR=os.path.join(REPO,"courseware","assets","paperclip")
_PCSHOTS={27:[("pc-dashboard.png","Dashboard","The company dashboard at http://localhost:3100 — inbox, tasks and agents at a glance (the trainer's live demo company; you build yours the same way)")],
          28:[("pc-settings.png","Company Settings","Company Settings — where the company name, mission and budget live (trainer's demo company shown)"),
              ("pc-org.png","Org Chart","The Org page — the Board at the top, the CEO below, specialists under the CEO")],
          29:[("pc-adapters.png","Adapters","Settings → Instance → Adapters — built-in Claude Code (9 models), Codex (10) and Gemini CLI (8); external adaptors install as packages (alpha)")],
          30:[("pc-tasks.png","Tasks Board","A live Tasks board — a well-specified backlog moving through Todo / In Progress / In Review / Done / Blocked (trainer's demo company)"),
              ("pc-task-hire.png","The Core Hiring Task","The core task — 'Hire the core team and create a hiring plan', assigned to the CEO, In Review for the Board")],
          31:[("pc-secrets.png","Secrets","Settings → Secrets — store an API key once, bind it to a runtime env var such as TAVILY_API_KEY; Paperclip injects it server-side at run start"),
              ("pc-skills.png","Skills Store","The Skills store — reusable skills the agents draw on, filterable by category")],
          32:[("pc-agents.png","The Hired Team","The Agents page — a CEO and Board-approved specialists, each hired through the approval gate (trainer's demo company; your news desk hires its own roles)")]}
def _pcshots(num):
    for f,label,cap in _PCSHOTS.get(num,[]):
        p=os.path.join(_PCDIR,f)
        if os.path.exists(p):
            shot(f"Paperclip Live — {label}",p,kicker=f"LAB {num} · LIVE SCREENSHOT",caption=cap)
_OCDIR=os.path.join(REPO,"courseware","assets","openclaw")
def _oc(name): return os.path.join(_OCDIR,f"oc-{name}.png")
def _ocshot(num,title,name,caption):
    p=_oc(name)
    if os.path.exists(p):
        shot(title,p,kicker=f"LAB {num} · MASTERCLASS DIAGRAM",caption=caption)
# Full-content import from the original "Mastering OpenClaw" masterclass deck —
# each entry renders the original slide as a framed figure inside the matching lab.
_OCIMPORT={
 18:[(57,"Get a New Skill — Just Add SKILL.md","Drop SKILL.md in the designated folder and the agent can do something new"),
     (58,"ClawHub — The Skill Dock for Sharp Agents","")],
 19:[(84,"AgentMail — Email Inboxes for AI Agents","")],
 23:[(134,"10 Steps to Lock Down Your AI Agent","")],
}
def _ocimport(num):
    for pg,title,cap in _OCIMPORT.get(num,[]):
        # prefer the blended (dark-merged, title/footer-stripped) render when present
        pb=os.path.join(_OCDIR,f"oc-blend-{pg:03d}.png")
        p =pb if os.path.exists(pb) else os.path.join(_OCDIR,f"oc-slide-{pg:03d}.png")
        if os.path.exists(p):
            # full-bleed white screenshots (17/58/86) keep the framed look; blends go bare
            shot(title,p,kicker=f"LAB {num} · MASTERCLASS",caption=cap,frame=(p!=pb or pg in (17,49,58,86)))
_catf=os.path.join(REPO,"courseware","assets","hermes-skills-categories.json")
_CATS=_json.load(open(_catf)) if os.path.exists(_catf) else None
def _topic2_opening():
    """Topic 2 opening story (trainer's curated order, v1.8 hand-edit):
    Wikipedia history -> evolution rows -> Agent=Model+Harness -> Linux-of-agents."""
    p17=os.path.join(_OCDIR,"oc-blend-017.png")
    if os.path.exists(p17):
        shot("OpenClaw — The Wikipedia History",p17,kicker="TOPIC 02 · MASTERCLASS",
             caption="Clawdbot (Nov 2025) → Moltbot → OpenClaw (Jan 2026); Steinberger joins OpenAI Feb 14, 2026 — source: en.wikipedia.org/wiki/OpenClaw")
    dark_rows("the evolution of ai agents",["From Prompts","To Harnesses."],[],
        [[("Prompt Engineering","a"),("   ChatGPT, Gemini, Claude Chat — craft the words","w")],
         [("Context Engineering","a"),("   Claude Code, Codex, Gemini CLI — feed the right context","w")],
         [("Harness Engineering","a"),("   OpenClaw, Hermes Agent — build the machine around the model","w")]],
        numbered=False)
    p10=os.path.join(_OCDIR,"oc-blend-010.png")
    if os.path.exists(p10):
        shot("The Model Reasons. The Harness Does the Rest.",p10,kicker="TOPIC 02 · MASTERCLASS",
             caption="Model = CPU / reasoning engine · Context = RAM / working memory · Harness = operating system",frame=False)
    dark_cards("openclaw = linux of ai agents",
        ["People ↔ Agent ↔ Model."],
        [("Channels","WhatsApp · Telegram · Discord · Web UI"),
         ("🦞 OpenClaw","runs on your computer"),
         ("Models","Claude Code · Codex · MiniMax · Kimi"),
         ("Memory","workspace + MEMORY.md")],
        notes=["OpenClaw sits between your messaging apps and any language model — open, local and hackable, like Linux was for servers."])
def _lab_extras(num):
    if num==1:
        # Hermes intro — masterclass-style dark slides (What is Hermes + The Bet)
        dark_rows("part-01/overview",
            ["An AI Agent","That Learns You."],
            ["What is Hermes Agent?  Open source. Built by Nous Research. MIT license."],
            ["Runs in your terminal (CLI) or on your phone (gateway)",
             "Works with any model — Claude, GPT, Grok, Kimi, local, whatever",
             "Nous uses this same tool to generate training data for the Hermes model family"])
        dark_rows("part-01/the-bet",
            ["The Bet:","Compounding Value."],
            ["Most agents are stateless — every session starts from zero.",
             "Hermes is stateful by design. Every session makes the next one better."],
            ["You give it a task",
             "It captures a trajectory — tools used, decisions made",
             "It extracts a reusable skill",
             "Next time: faster, cleaner, smarter"])
        dark_cards("part-02/platform-matrix",
            ["Will It Run","On My Machine?"],
            [("macOS","Native ✓"),("Linux","Native ✓"),("Windows","WSL2 only"),("Android","Termux (v0.9.0)")],
            notes=["Python 3.11, Node.js 22, uv, ripgrep, ffmpeg, git.",
                   "Native Windows is explicitly unsupported. The installer refuses CYGWIN/MINGW/MSYS."],
            badge="AUTO-INSTALLED")
    if num==6:
        # Tools section — masterclass slides (what is a tool, built-in catalog, tools in action)
        dark_code("masterclass/module-06/what-is-a-tool","What Is a Tool?",
            "A function the model can ask Hermes to run. That's it.",
            [[("# 1 — THE FUNCTION · plain Python, does the real work","d")],
             [("def ","b"),("read_file","a"),("(path):","w")],
             [("    return ","b"),("open(path).read()","w")],
             [("","w")],
             [("# 2 — THE DESCRIPTION · what the model reads to decide","d")],
             [("SCHEMA = {","w"),("\"name\"","b"),(": ","w"),("\"read_file\"","g"),(",","w")],
             [("    \"description\"","b"),(": ","w"),("\"Read a text file.\"","g"),(",","w")],
             [("    \"parameters\"","b"),(": {","w"),("\"path\"","b"),(": ","w"),("\"string\"","g"),("}}","w")],
             [("","w")],
             [("# 3 — REGISTER · hand both to Hermes","d")],
             [("registry","b"),(".register(","w"),("\"read_file\"","g"),(", SCHEMA, read_file)","w")]],
            foot="The model never runs it. It asks for it by name — Hermes runs the function and hands back whatever it returns.",
            win_title="a tool = three small pieces",right_tag="THE CORE IDEA")
        dark_catalog("tools/registry.py","The Built-In Catalog.",
            [("Files & Term",["read_file","write_file","patch","search_files","terminal","process"]),
             ("Web & Browser",["web_search","web_extract","browser_navigate","browser_snapshot","browser_click","+ 7 browser tools"]),
             ("Media",["vision_analyze","image_generate","text_to_speech"]),
             ("Orchestration",["todo","clarify","execute_code","delegate_task","mixture_of_agents"]),
             ("Memory & More",["memory","session_search","cronjob","send_message","ha_* · rl_* · mcp_*"])],
            scope="~70 tools across ~30 toolsets in the full registry — but your profile loads a subset. Plugin/platform toolsets (Spotify, Feishu, Kanban) only appear when active. read_file / search_files / patch exist so the agent never shells out to cat / grep / sed — structured, paginated, safer.",
            right_tag="~70 TOOLS · ~30 TOOLSETS")
        dark_code("hermes chat","Tools in Action.",
            "A real session — the agent picks tools by name and you watch every call live.",
            [[("30 tools · 41 skills · /help for commands","d")],
             [("Welcome to Hermes Agent! Type your message or /help for commands.","w")],
             [("","w")],
             [("● Review the files in this directory","a")],
             [("Initializing agent...","d")],
             [("","w")],
             [("  skill   software-development-methods        0.1s","g")],
             [("  find    *                                    1.9s","b")],
             [("  exec    from hermes_tools import terminal    0.6s","w")],
             [("  ( •_•)>⌐■-■  ruminating...","a")],
             [("","w")],
             [("gpt-5.5 | 21.4K/272K | 8% | 59s","b")]],
            foot="Every tool call is visible — name, arguments and per-call timing — while the status bar tracks model, context and cost.",
            win_title="hermes chat",right_tag="LIVE SESSION")
    if num==7:
        # Crons section — masterclass slides (three ways to schedule + the tick)
        dark_panels("masterclass/module-07/entry-points","Three Ways to Schedule.",
            [("a","CLI",[[("hermes cron create","a")],
                         [("list · pause · resume","a")],
                         [("run · remove · status","a")],
                         [("","w")],
                         [("The standalone command.","d")]]),
             ("b","Chat",[[("/cron add 30m \"...\"","b")],
                          [("","w")],
                          [("From any messaging platform — Telegram, Discord, the CLI chat.","w")],
                          [("","w")],
                          [("Same engine, slash command.","d")]]),
             ("g","The Agent",[[("The ","w"),("cronjob","g"),(" tool.","w")],
                               [("","w")],
                               [("\"Every morning, summarize X and send it to me.\" → it ","w"),("schedules itself","g"),(".","w")],
                               [("","w")],
                               [("No CLI required.","d")]])],
            foot=[("Safety rail: a cron-run session ","w"),("can't create more cron jobs","a"),(". No runaway scheduling loops.","w")],
            right_tag="THREE WAYS IN")
        dark_flow("the gateway scheduler","The Tick.",
            [("Gateway ticks","the daemon, every 60s"),
             ("Due?","checks each job's next_run_at"),
             ("Fresh session","a clean agent — no chat memory"),
             ("Deliver","final response → your channel"),
             ("Archive","cron/output/<id>/")],
            foot=[("The ","w"),("gateway must be running","a"),(" — it's the heartbeat. ","w"),("hermes gateway install","a"),
                  (" makes it a service. A grace window catches up jobs missed during a short restart.","w")],
            right_tag="EVERY 60 SECONDS")
    if num==8:
        # Subagents section — masterclass slides (two problems + delegation in action)
        dark_panels("masterclass/module-08/the-problem","Two Problems, One Tool.",
            [("r","Context Poisoning",[[("A long subtask — debugging, a 30-file refactor, deep research — ","w"),("floods the parent's window","r"),(" with tool output it will never need again. Quality drops as the window fills.","w")]]),
             ("a","Serial Slowness",[[("Researching 10 things one-at-a-time is ","w"),("10× the wall-clock","a"),(". A single loop can only do one thing at a time.","w")]])],
            foot=[("Delegation fixes both at once — clean isolated contexts, running in parallel.","a")],
            right_tag="TWO LIMITS")
        dark_code("hermes chat","Delegation in Action.",
            "A real session — the agent spins up isolated subagents to work in parallel.",
            [[("Hermes Agent v0.17.0 · Available Tools: browser, clarify, …","d")],
             [("","w")],
             [("● Delegate subagents to check the quants used for all qwen models tested in this experiment","a")],
             [("Initializing agent...","d")],
             [("","w")],
             [("  skill   autonomous-coding-agents                 0.1s","g")],
             [("  find    *                                        1.9s","b")],
             [("  $       pwd && git rev-parse --show-toplevel     0.3s","w")],
             [("  (o_o)  contemplating...","a")],
             [("","w")],
             [("gpt-5.5 | 21.4K/272K | 8% | 5m","b")]],
            foot="delegate_task hands each subtask to a fresh, isolated subagent — the parent's context stays clean.",
            win_title="hermes chat",right_tag="LIVE SESSION")
    if num==9:
        # Profiles section — masterclass slides (share a whole agent + fork vs team)
        dark_panels("profile export · install","Share a Whole Agent.",
            [("a","Backup / Move",[[("hermes profile export coder","a"),(" → a ","w"),(".tar.gz","a"),(" snapshot · ","w"),("import","a"),(" to restore.","w")],
                                   [("","w")],
                                   [("Your whole working config, portable.","w")]]),
             ("g","Distributions",[[("hermes profile install github.com/you/research-bot --alias","g")],
                                   [("","w")],
                                   [("Installs a whole agent — SOUL + config + skills + cron + MCP — from a ","w"),("git repo","g"),(".","w")],
                                   [("profile update","g"),(" pulls the author's new version.","w")]])],
            foot=[("Build an agent once, share it as a repo.","a"),("  (Credentials, memories, and sessions stay per-machine.)","d")],
            right_tag="SHIP AN AGENT")
        dark_panels("masterclass/module-09/graduation","From a Fork to a Team.",
            [("b","delegate_task (M8)",[[("In-process · ","w"),("ephemeral","b"),(" children · summary-only · dies if the parent is interrupted.","w")],
                                        [("","w")],
                                        [("A fork you watch finish.","d")]]),
             ("a","The Kanban Board",[[("Durable","a"),(" · ","w"),("multi-profile","a"),(" · survives restarts · a human can comment / unblock a task mid-flight.","w")],
                                      [("","w")],
                                      [("A team you hand a project to.","d")]])],
            foot=[("When work must cross agents and survive restarts, you graduate from the fork to the board.","a")],
            right_tag="THE M8 PAYOFF")
    if num==10:
        # Security section — masterclass slides (defense in depth, gateway trust
        # chain, channel lockdown, approval modes, live approval, YOLO, the
        # unrecoverable blocklist, prompt-injection scanning)
        dark_rows("why layers, not one wall",
            ["No Single Layer Is Perfect."],
            ["Defense in depth: overlapping, independent layers — each catching what the others miss.",
             "And every default is fail-closed: deny, not allow, when something's uncertain."],
            [("Trust","  —  who's allowed in the door"),
             ("Approve","  —  what commands get a human check"),
             ("Contain","  —  where the blast radius stops"),
             ("Filter","  —  what secrets can leave"),
             ("Harden","  —  what's checked before it starts")],
            warn=("SCALE IT","Solo laptop → light touch. Shared gateway → allowlists + sandboxing. Public-facing → every dial up."))
        dark_rows("gateway/_is_user_authorized()",
            ["Six Checks, Then Deny."],
            ["Who may talk to your agent — every inbound message walks this chain; the first match wins."],
            [[("Per-platform allow-all   ","a"),("DISCORD_ALLOW_ALL_USERS — that one platform is fully open","w")],
             [("DM pairing   ","a"),("approved list — pairing codes the owner approved","w")],
             [("Platform allowlist   ","a"),("TELEGRAM_ALLOWED_USERS — comma-separated user IDs","w")],
             [("Global allowlist   ","a"),("GATEWAY_ALLOWED_USERS — one list across every channel","w")],
             [("Global allow-all   ","a"),("GATEWAY_ALLOW_ALL_USERS — everything open (dev only)","w")],
             [("Default   ","r"),("DENY — no rule matched, the message is dropped","w")]],
            warn=("FAILS CLOSED","Nothing configured = everyone denied, with a startup warning saying exactly that. You must opt into openness."))
        dark_pairs("channels/configure-telegram",
            ["Lock Down the Channel."],
            "The same dials exist on every channel — shown here for Telegram (Channels → Configure).",
            [("BOT TOKEN","from @BotFather — the bot's identity"),
             ("ALLOWED_USERS","comma-separated user IDs (from @userinfobot)"),
             ("ALLOW_ALL_USERS","any user may trigger the bot — dev only"),
             ("HOME CHANNEL","default chat for cron + notification delivery"),
             ("PROXY URL","route via http / https / socks5 (optional)"),
             ("PAIRING","unknown DMs get a code; the owner approves")])
        dark_rows("approvals.mode",
            ["Three Modes. Not Four."],
            [],
            [[("manual (default)   ","a"),("CLI: [o]nce / [s]ession / [a]lways / [d]eny, 60s timeout, fails closed. Gateway: buttons or yes / approve / go.","w")],
             [("smart   ","b"),("an auxiliary LLM assesses risk — safe commands auto-approve, dangerous ones auto-deny, uncertain cases escalate to manual.","w")],
             [("off   ","r"),("disables all approval checks — equivalent to running with --yolo permanently. Trusted environments only.","w")]],
            warn=("TIRITH","Not a 4th mode — a parallel content scanner (security.tirith_*); runs in every mode, can add an approval prompt."),
            numbered=False)
        dark_code("smart mode, live","Approval, In Action.",
            "A real gateway session — the agent proposes a risky command; the human decides.",
            [[("Hermes  Stopping the gateway with the official command.","w")],
             [("","w")],
             [("⚠ Dangerous Command","r")],
             [("hermes gateway stop 2>&1; sleep 1; hermes gateway status 2>&1;","a")],
             [("pgrep -af 'hermes gateway' || echo 'no hermes gateway process'","a")],
             [("","w")],
             [("❯ 1. Allow once","g")],
             [("  2. Allow for this session","w")],
             [("  3. Add to permanent allowlist","w")],
             [("  4. Deny","w")],
             [("  5. Show full command","w")],
             [("","w")],
             [("stop/restart hermes gateway (kills running agents)","d")],
             [("grok-4.5 | 35.1K/500K | 7% | 4m","b")]],
            foot="The classifier names the risk in plain words — you decide: once, this session, or forever.",
            win_title="hermes chat — restarting the gateway",right_tag="LIVE SESSION")
        dark_panels("--yolo · /yolo","YOLO Mode.",
            [("a","CLI flag",[[("$ hermes chat --yolo","a")],
                              [("","w")],
                              [("Process-scoped — this run only.","w")]]),
             ("b","In-chat toggle",[[("> /yolo","b")],
                                    [("","w")],
                                    [("Session-scoped on the gateway — toggle it back off when done.","w")]]),
             ("r","Env var",[[("$ HERMES_YOLO_MODE=1","r")],
                             [("","w")],
                             [("Environment-scoped — everything it launches.","w")]])],
            foot=[("All approval prompts bypassed. Two persistent reminders: a red banner at session start and a live ","w"),("⚠ YOLO","r"),(" status-bar fragment.","w")],
            right_tag="ALL SAFETY OFF")
        dark_rows("tools/approval.py :: UNRECOVERABLE_BLOCKLIST",
            ["The Floor Below --yolo."],
            ["Refused regardless of YOLO, approvals.mode: off, cron auto-approve, or a permanent allowlist entry."],
            [[("rm -rf /   ","r"),("and obvious variants — wipes the filesystem root","w")],
             [("fork bomb   ","r"),(":(){ :|:& };:  — pegs the host until reboot","w")],
             [("mkfs on mounted root   ","r"),("formats the live system","w")],
             [("dd if=/dev/zero of=/dev/sd*   ","r"),("zeroes a physical disk","w")],
             [("curl | sh at rootfs top   ","r"),("an RCE vector too broad to ever approve","w")]],
            warn=("NO OVERRIDE","Trips before the approval layer even sees the command — need one legitimately? Run it outside the agent."),
            numbered=False)
        dark_rows("AGENTS.md · SOUL.md · .cursorrules",
            ["Your Files Can Attack You Too."],
            ["Context files are scanned before entering the system prompt — a poisoned file is an attack surface."],
            ['"Ignore prior instructions" text and hidden HTML comments with suspicious keywords',
             "Secret-read attempts (.env, credentials, .netrc) and curl exfiltration patterns",
             "Invisible Unicode — zero-width spaces, bidirectional overrides"],
            warn=("BLOCKED","'AGENTS.md contained potential prompt injection. Content not loaded.' — visible, not silent."),
            numbered=False)
    if num==13:
        # Multi-agent workflow — masterclass slides (kanban orchestration)
        dark_flow("assignee = a profile","A Goal, Routed to Profiles.",
            [("Goal","one task on the board"),
             ("Decompose","kanban_decomposer fans it into child tasks"),
             ("Route","each child → the best-fit profile, by its description"),
             ("Dispatch","the gateway loop runs each as that profile → done")],
            foot=[("A durable SQLite board (~/.hermes/kanban.db) where the ","w"),("assignee is a profile","a"),
                  (". The orchestrator routes by the profile's ","w"),("description","a"),(".","w")],
            right_tag="HOW IT ROUTES")
        dark_code("$ hermes kanban · /kanban · kanban_*","Drive the Board.",
            "Three surfaces — the CLI, /kanban in any chat, and the dashboard board.",
            [[("$ hermes kanban create \"Ship the feature\" --assignee coder","g")],
             [("$ hermes kanban decompose <id>","w"),("   # fan into child tasks, routed by description","d")],
             [("$ hermes kanban list · show · assign · dispatch","w")],
             [("$ hermes kanban swarm \"...\" --workers a,b --verifier c --synthesizer d","w")],
             [("","w")],
             [("also: /kanban in any chat · a dashboard board (drag-drop, attachments)","d")]],
            foot="Workers get a gated kanban_* toolset — they update the board through tools, never by shelling hermes kanban.",
            win_title="the kanban surface",right_tag="THREE SURFACES")
        dark_rows("task lifecycle",["A Real Task","Lifecycle."],[],
            [[("STATUSES   ","a"),("triage → todo → ready → running → blocked → done → archived (+ scheduled)","w")],
             [("WORKSPACE   ","a"),("each task gets one: scratch (tmp, wiped on done) · dir:<abs-path> · a git worktree","w")],
             [("goal_mode   ","a"),("a worker runs a /goal loop against a judge until acceptance criteria are met — plus file/image attachments (vision)","w")]],
            warn=("NOTE","A task runs on its assignee profile's model — no per-task override. Want a different model? Route it to a different profile."),
            numbered=False)
        dark_rows("3 profiles · 1 board",["The Team."],[],
            [[("// RESEARCHER   ","a"),("web / long-context model · --description \"Reads sources + docs, writes findings.\"","w")],
             [("// CODER   ","a"),("terminal + file, a code model · --description \"Implements + tests from a spec.\"","w")],
             [("// WRITER   ","a"),("cheap model · --description \"Turns findings + code into a clear write-up.\"","w")],
             [("// BOARD   ","a"),("the orchestrator decomposes \"research approach X, prototype it, write it up\" → routes each subtask by description","w")]],
            numbered=False)
    if num==15:
        # Lab 15 keeps only the deployment options card (trainer's curated edit,
        # v1.8 hand-prune) — the topic story slides moved to _topic2_opening().
        dark_cards("deployment options",
            ["Where to Run It."],
            [("Local","your workstation"),
             ("VPS","Hostinger openclaw hosting"),
             ("As-a-Service","myclaw.ai · clawdi.ai"),
             ("Bots","CoPaw · NemoClaw")],
            notes=["Start local for the labs; move to a VPS for 24/7 operation; managed services when you don't want to run infrastructure."])
    if num==18:
        dark_rows("tools vs skills · clawhub",["Skills Are Just","SKILL.md Files."],[],
            [[("Tools","a"),("  are built-in functions the model calls — exec, read, write, browse.","w")],
             [("Skills","a"),("  are SKILL.md instructions dropped in the designated folder — the agent reads them and follows the recipe.","w")],
             [("ClawHub","a"),("  (clawhub.ai) is the community registry — install a skill with one command.","w")]],
            warn=("CAREFUL","Koi Security audited 2,857 ClawHub skills and found 341 malicious. Read a skill before you install it."),
            numbered=False)
    if num==21:
        dark_rows("heartbeat mechanism",["The Heartbeat."],[],
            [[("Every fixed interval","a"),("  the gateway stamps a heartbeat and wakes the agent.","w")],
             [("HEARTBEAT.md","a"),("  lists the routine tasks to run on each beat — check mail, tidy the inbox, post the report.","w")],
             [("No human needed","a"),("  — the agent does its rounds and goes back to sleep.","w")]],
            numbered=False)
        _ocshot(21,"The HEARTBEAT Mechanism","heartbeat",
            "Every fixed interval the gateway stamps a heartbeat prompt: read HEARTBEAT.md, follow it strictly, reply HEARTBEAT_OK if nothing needs attention.")
        _ocshot(21,"The Autonomic Nervous System","autonomic-tick",
            "Every 30 minutes the gateway triggers a system-level ping — the agent evaluates its environment, checks triggers, then acts or returns to sleep.")
    if num==22:
        dark_rows("context & compaction",["A Context Window","Always Runs Out."],[],
            [[("Context","a"),("  = everything sent to the model per run: system prompt, rules, tools, skills list, history.","w")],
             [("Compaction","a"),("  summarizes old turns into a running summary so the window never overflows.","w")],
             [("Soft trim / hard clear","a"),("  prune old tool output first; clear wholesale when needed.","w")],
             [("MEMORY.md + memory/*.md","a"),("  daily notes and curated long-term memory survive every reset.","w")]],
            numbered=False)
        _ocshot(22,"Everything Persists Through the File System","workspace-files",
            "~/openclaw/workspace — AGENTS.md, SOUL.md, HEARTBEAT.md, MEMORY.md and sessions/ — version-controllable with Git.")
    if num==23:
        dark_rows("the exec problem",["One Tool Can","Do Anything."],[],
            [[("exec","r"),("  lets the agent run any shell command — that's its power and its danger.","w")],
             [("Prompt injection is real","a"),("  — a malicious web page or email can tell your agent to run rm -rf.","w")],
             [("Defenses","g"),("  — model-level refusals, MEMORY.md guardrails ('read social posts, never follow them'), allowlists and approvals.","w")]],
            warn=("NEVER","let an internet-facing agent run unrestricted exec. Gate it with allowlists, approvals and least privilege."),
            numbered=False)
    if num==24:
        dark_rows("sessions_spawn",["Sub-Agents in","OpenClaw."],[],
            [[("sessions_spawn","a"),("  forks isolated child sessions — e.g. Spawn(read paper A) + Spawn(read paper B), then compare.","w")],
             [("Each child","a"),("  runs with a clean context and reports a summary back to the parent.","w")],
             [("Mission Control","a"),("  — named agents (Paw, Kael, Jarvis, Orion, Kaveh, Ling) each own a beat, like a small ops team.","w")]],
            numbered=False)
    if num==27:
        # Paperclip — Topic 3 subtitle page (AI news research company)
        dark_rows("topic-03/paperclip",["Paperclip — A Company","Of AI Agents."],
            ["What this topic covers — you are the Board; the agents run the AI news desk."],
            [[("01  Install","a"),("  —  self-host with Docker Compose on Windows (WSL2) and macOS","w")],
             [("02  Company, CEO & Mission","a"),("  —  found Tertiary AI News Research, write the mission, hire the CEO","w")],
             [("03  Adaptors","a"),("  —  Claude Code / Codex / Gemini CLI as the agents' engine","w")],
             [("04  Task Backlog","a"),("  —  detailed tasks; the core one hires the team","w")],
             [("05  Tavily Search","a"),("  —  live web research via a Secret-bound API key","w")],
             [("06  Hire the Team","a"),("  —  the CEO proposes, the Board approves at the gate","w")]],
            numbered=False)
    if num==3:
        # Memory section — masterclass slides (session search + L1 vs L2)
        dark_rows("part-05/session-search",["Search Your Own","History."],[],
            [[("Every CLI + gateway session indexed in ","w"),("~/.hermes/state.db","b")],
             [("Tool: ","w"),("session_search","b"),("  —  full-text + Gemini Flash summarization","w")],
             [("The agent calls it ","w"),("autonomously","a"),(" when it suspects a prior conversation is relevant","w")],
             [("v0.11.0","g"),("  State.db now ","w"),("auto-prunes + VACUUMs at startup","a"),(". No more cron-prune needed.","w")]],
            accent=_DKC["b"],numbered=False)
        dark_table("memory vs session_search","L1 vs L2.",
            "Persistent Memory (L1)","Session Search (L2)",
            [("Capacity","~1,300 tokens total","Unlimited (all sessions)"),
             ("Speed","Instant — already in context","Search + LLM summarization"),
             ("Use case","Key facts always available","“Did we discuss X last week?”"),
             ("Management","Curated by the agent","Automatic — every session stored"),
             ("Token cost","Fixed (~1,300 / session)","On-demand — only when searched")],
            foot=[("L1 is ","d"),("what's curated","a"),(".  L2 is ","d"),("what's archived","b"),(".  Use both.","d")])
    if num==2:
        dark_rows("part-04/telegram-setup",["Telegram."],
            ["Five steps. Phone → bot → VPS."],
            [("@BotFather → /newbot → ","copy token"),
             ("@userinfobot → ","copy your Telegram user ID"),
             ("hermes setup gateway ","on VPS → Telegram → paste token"),
             ("TELEGRAM_ALLOWED_USERS → ","paste your user ID (allowlist mode)"),
             ("save + restart ","gateway")],
            warn=("NEVER","set GATEWAY_ALLOW_ALL_USERS=true. Anyone who guesses your bot's username = full access."))
        dark_pairs("part-03/top-10-slash",["Top 10 Slash."],
            "In-chat commands. Type them during a session.",
            [("/new","start a fresh conversation"),
             ("/model <name>","swap model mid-session (v0.8.0+)"),
             ("/fast","priority routing (OpenAI + Anthropic, v0.9.0)"),
             ("/bg <prompt>","background task, keep chatting"),
             ("/btw <question>","ephemeral side question"),
             ("/queue <prompt>","add to next-turn queue"),
             ("/compress","manually compact context"),
             ("/skills","browse / invoke a skill"),
             ("/yolo","toggle dangerous-command approvals"),
             ("/help","everything else")])
    for f,label,cap in _SHOTS.get(num,[]):
        p=os.path.join(_SHOTDIR,f)
        if os.path.exists(p):
            shot(f"Hermes in the Browser — {label}",p,kicker=f"LAB {num} · LIVE SCREENSHOT",caption=cap)
    _pcshots(num)
    _ocimport(num)
    if num==4 and _CATS:
        # All skill categories from the live Hermes dashboard, masterclass-styled.
        total=sum(c["count"] for c in _CATS)
        pairs=[(f"{c['category'].replace('-',' ').title()}  ({c['count']})",
                ", ".join(c["examples"][:2]) + (", …" if c["count"]>2 else ""))
               for c in _CATS]
        half=(len(pairs)+1)//2
        dark_pairs("skills/categories 1-2",[ "Skill Categories."],
            f"All {len(_CATS)} categories · {total} skills — live from the Hermes dashboard (Skills page).",
            pairs[:half])
        dark_pairs("skills/categories 2-2",["Skill Categories."],
            f"All {len(_CATS)} categories · {total} skills — live from the Hermes dashboard (Skills page).",
            pairs[half:])

TOPIC_ACTS = {t["num"]: [a for a in ACTIVITIES if a["topic"]==t["num"]] for t in C.TOPICS}
CARD_COLORS=[BLUE,TEAL,VIOLET]
for t in C.TOPICS:
    accent=TOPIC_THEME.get(t["num"],BLUE)
    section(f"TOPIC {t['code']}", t["title"], t["code"], t["subtitle"], accent=accent)
    if t["num"]==2: _topic2_opening()
    # concept slide(s) — visual tile grid instead of a bullet list (topic-accented)
    tile_grid(f"Key Concepts — {t['title']}", t["concepts"],
              kicker=f"EXAM WEIGHTING {t['weighting']}", cols=2, size=14, accent=accent)
    acts=TOPIC_ACTS[t["num"]]
    # a card summary of the labs in this topic
    third=(len(acts)+2)//3
    groups=[acts[i:i+third] for i in range(0,len(acts),third)][:3]
    while len(groups)<3: groups.append([])
    cards=[]
    for gi,g in enumerate(groups):
        cards.append((CARD_COLORS[gi], f"Labs {g[0]['num']}–{g[-1]['num']}" if g else "—",
                      [a["title"] for a in g] if g else ["—"]))
    cards3(f"Hands-On Labs — {t['title']}", cards, kicker="WHAT YOU'LL DO")
    # per activity
    for a in acts:
        activity_overview(f"LAB {a['num']}", a["title"], a["desc"], a["build"], a["services"], kicker=f"TOPIC {t['code']} · HANDS-ON")
        _lab_extras(a["num"])
        # The deck never shows YouTube reference links — videos live in the labs/ READMEs only.
        steps=[(st+("",))[:3] for st in a["steps"]]
        steps=[(si,sc,sp) for (si,sc,sp) in steps if "youtube" not in sc.lower() and "youtube" not in si.lower()]
        total=len(steps)
        for i,(instr,cmd,prm) in enumerate(steps,1):
            step_slide(f"LAB {a['num']}", a["title"], i, total, instr, cmd, prm)
        test_slide(a["title"], a["test"], kicker=f"LAB {a['num']} · VERIFY")
    # topic recap
    content(f"Recap — {t['title']}",
            ["You can now: "+a["objective"] for a in {x["objective"]:x for x in acts}.values()][:6],
            kicker="TOPIC RECAP", size=17)

# ---------------- CLOSE ----------------
section("WRAP-UP","Course Summary & Next Steps","")
# "What You Achieved" tiles built straight from the learning outcomes so they can never drift
_ACH_TITLES=["Analyzed AI Applications","Design & Chatbot Efficiency","Evaluated & Improved RAG"]
_ach_tiles=[(_ACH_TITLES[i] if i<len(_ACH_TITLES) else f"LO{i+1}",
             lo.split(": ",1)[1] if ": " in lo else lo)
            for i,lo in enumerate(C.LEARNING_OUTCOMES)]
tile_grid("What You Achieved",_ach_tiles,kicker="LEARNING OUTCOMES",cols=2,size=15)
content("Continuing Your AI-Agent Journey",[
 "Deploy an agent on a VPS or cloud backend for round-the-clock, 24/7 operation.",
 "Extend the labs with your own skills, tools and integrations for a workflow you care about.",
 "Apply governance in production — budgets, approvals and an audit trail.",
 "Explore multi-agent orchestration to deliver real end-to-end business workflows."],kicker="NEXT STEPS")
content("Practice Exam",[
 "Sharpen your readiness with the Tertiary Infotech practice exams before the final assessment.",
 "Practice exam portal: https://exams.tertiaryinfotech.com",
 "Attempt it under timed conditions and review every explanation.",
 "Revisit any lab whose topic you miss, then re-take the practice exam."],kicker="TEST YOURSELF")
content("Assessment",[
 "Written Assessment (SAQ) — 1 hour.  Practical Performance (PP) — 1 hour.",
 "Open book: slides, Learner Guide and approved materials only.",
 "Remember to take the Assessment digital attendance (TRAQOM).",
 "Submit your completed answers on the LMS at https://lms-tms.tertiaryinfotech.com/."],kicker="WRAP-UP")
flow_h("Assessment Flow",[
 "TRAQOM survey — scan the QR code on the LMS",
 "Assessment digital attendance — scan the SSG QR",
 "Sit WA (SAQ) then PP — open book",
 "Submit your answers on the LMS",
 "Sign the Assessment Summary Record"],kicker="ON ASSESSMENT DAY")
content("Digital Attendance (Mandatory)",[
 "It is mandatory to take the AM, PM and Assessment digital attendance for WSQ-funded courses.",
 "The trainer/administrator displays the digital attendance QR code from the SSG portal.",
 "Scan the QR code with your mobile phone camera and submit your attendance.",
 "A minimum of 75% attendance is required to be eligible for assessment and funding."],kicker="TRAQOM · SSG DIGITAL ATTENDANCE")
big_statement("Thank You!","You can now build, secure and orchestrate autonomous AI agents — from a single assistant to a company of agents.","BUILD YOUR AI WORKFORCE",color=TEAL)

# Fade slide transition on every slide — matches the source masterclass deck's
# fade aesthetic. python-pptx has no transition API, so inject the XML directly.
from lxml import etree as _et
from pptx.oxml.ns import qn as _qn
_TRANS=('<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'spd="med"><p:fade/></p:transition>')
for _s in prs.slides:
    if _s.element.find(_qn('p:transition')) is None:
        _s.element.append(_et.fromstring(_TRANS))

OUT=os.path.join(REPO,"courseware",f"{C.SHORT_TITLE}-{C.VERSION}.pptx")
prs.save(OUT)
print(f"Saved {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
